from PyPDF2 import PdfReader, PdfWriter
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches, Pt
import pdfplumber
from pdf2image import convert_from_path
import os

# New LangChain imports for content chunking
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter

def split_pdf(pdf_path, output_dir):
    """
    Splits a PDF into individual page files and saves them.

    Args:
        pdf_path (str): Path to the input PDF.
        output_dir (str): Directory where split pages will be saved.

    Returns:
        int: Total number of pages split.
    """
    reader = PdfReader(pdf_path)
    for i, page in enumerate(reader.pages):
        writer = PdfWriter()
        writer.add_page(page)
        output_path = f"{output_dir}/page_{i+1}.pdf"
        with open(output_path, "wb") as f:
            writer.write(f)
    return len(reader.pages)

def convert_pdf_to_word(pdf_path, word_path):
    """
    Converts a PDF file to a Word (.docx) document.
    """
    cv = Converter(pdf_path)
    cv.convert(word_path, start=0, end=None)
    cv.close()

def convert_pdf_to_pptx(pdf_path, pptx_path):
    """
    Converts a PDF to a PowerPoint presentation with each page as a slide.
    """
    presentation = Presentation()
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)

    images = convert_from_path(pdf_path, dpi=150)

    for i, image in enumerate(images):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Blank layout
        temp_img = f"temp_page_{i}.jpg"
        image.save(temp_img)

        slide.shapes.add_picture(
            temp_img,
            left=Inches(0),
            top=Inches(0),
            width=presentation.slide_width,
            height=presentation.slide_height
        )

        os.remove(temp_img)

    presentation.save(pptx_path)
    return pptx_path

def get_split_chunks(pdf_path, chunk_size=1000, chunk_overlap=150):
    """
    Loads a PDF and returns its text content split into chunks for processing.

    Args:
        pdf_path (str): Path to the input PDF.
        chunk_size (int): Number of characters per chunk.
        chunk_overlap (int): Overlap between chunks.

    Returns:
        List[Document]: List of LangChain Document chunks with metadata.
    """
    loader = PyPDFLoader(pdf_path)
    docs = loader.load()
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return splitter.split_documents(docs)